#!/usr/bin/env python3
"""PPT Processor Tool - 提供 PPTX 处理相关功能

功能:
- 扫描 PPTX 文件
- 抽取页面内容（结构化 JSON + Markdown）
- 生成高清页面图片（通过 LibreOffice + Poppler）
- 保存讲稿
- 生成 TTS 音频
- 生成视频片段并合成

依赖:
- python-pptx: PPTX 解析
- LibreOffice: PPTX -> PDF 转换（系统依赖）
- Poppler (pdftoppm): PDF -> PNG 转换（系统依赖）
- ffmpeg: 音视频处理（系统依赖）
- Azure Speech SDK: TTS 生成
"""

from __future__ import annotations

import json
import os
import re
import shutil
import subprocess
import tempfile
import time
from concurrent.futures import ThreadPoolExecutor, as_completed
from pathlib import Path
from typing import Any
from xml.sax.saxutils import escape as _xml_escape


def _ensure_dir(path: Path) -> Path:
    """确保目录存在"""
    path.mkdir(parents=True, exist_ok=True)
    return path


def _safe_filename(name: str, max_len: int = 80) -> str:
    """生成安全的文件名"""
    s = (name or "").strip()
    s = re.sub(r"\s+", "_", s)
    s = re.sub(r"[^A-Za-z0-9_\-\u4e00-\u9fff]+", "_", s)
    s = s.strip("_-")
    return (s[:max_len] or "pptx")


def scan_pptx_files(input_path: str, output_dir: str) -> dict[str, Any]:
    """扫描输入路径，检测所有 PPTX 文件
    
    Args:
        input_path: PPTX 文件路径或包含 PPTX 的目录
        output_dir: 输出目录
    
    Returns:
        检测到的 PPTX 文件列表
    """
    input_p = Path(input_path).expanduser().resolve()
    output_p = Path(output_dir).expanduser().resolve()
    _ensure_dir(output_p)
    
    pptx_files: list[dict[str, str]] = []
    
    if input_p.is_file():
        if input_p.suffix.lower() in ('.pptx', '.ppt'):
            pptx_files.append({
                "path": str(input_p),
                "stem": input_p.stem,
                "name": input_p.name
            })
    elif input_p.is_dir():
        for f in sorted(input_p.glob("**/*.pptx")):
            pptx_files.append({
                "path": str(f),
                "stem": f.stem,
                "name": f.name
            })
        for f in sorted(input_p.glob("**/*.ppt")):
            pptx_files.append({
                "path": str(f),
                "stem": f.stem,
                "name": f.name
            })
    else:
        return {"ok": False, "error": f"输入路径不存在: {input_path}"}
    
    if not pptx_files:
        return {"ok": False, "error": f"未找到 PPTX 文件: {input_path}"}
    
    # 保存扫描结果
    result_path = output_p / "pptx_files.json"
    result_path.write_text(json.dumps(pptx_files, ensure_ascii=False, indent=2), encoding="utf-8")
    
    return {
        "ok": True,
        "count": len(pptx_files),
        "files": pptx_files,
        "result_file": str(result_path)
    }


def _extract_text_from_shape(shape) -> str:
    """从形状中提取文本"""
    texts = []
    if hasattr(shape, "text") and shape.text:
        texts.append(shape.text.strip())
    if hasattr(shape, "text_frame"):
        for paragraph in shape.text_frame.paragraphs:
            para_text = "".join(run.text for run in paragraph.runs).strip()
            if para_text:
                texts.append(para_text)
    return "\n".join(texts)


def _extract_slide_content(slide, slide_idx: int) -> dict[str, Any]:
    """提取单页幻灯片内容"""
    content = {
        "page_number": slide_idx + 1,
        "title": "",
        "texts": [],
        "notes": "",
        "shapes": []
    }
    
    # 提取标题
    if slide.shapes.title:
        content["title"] = slide.shapes.title.text.strip()
    
    # 提取所有形状中的文本
    for shape in slide.shapes:
        shape_info = {
            "type": shape.shape_type.name if hasattr(shape.shape_type, "name") else str(shape.shape_type),
            "text": ""
        }
        
        text = _extract_text_from_shape(shape)
        if text:
            shape_info["text"] = text
            if text not in content["texts"] and text != content["title"]:
                content["texts"].append(text)
        
        content["shapes"].append(shape_info)
    
    # 提取备注
    if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
        notes_text = slide.notes_slide.notes_text_frame.text.strip()
        if notes_text:
            content["notes"] = notes_text
    
    return content


def _convert_pptx_to_images(pptx_path: Path, output_dir: Path, dpi: int = 150) -> list[str]:
    """将 PPTX 转换为高清图片
    
    使用 LibreOffice 转 PDF，再用 Poppler 转 PNG
    """
    images = []
    images_dir = _ensure_dir(output_dir / "images")
    
    # 创建临时目录用于 PDF
    with tempfile.TemporaryDirectory() as temp_dir:
        temp_path = Path(temp_dir)
        
        # 使用 LibreOffice 转换 PPTX 到 PDF
        pdf_path = temp_path / f"{pptx_path.stem}.pdf"
        
        try:
            # LibreOffice 命令
            lo_cmd = [
                "libreoffice",
                "--headless",
                "--convert-to", "pdf",
                "--outdir", str(temp_path),
                str(pptx_path)
            ]
            
            result = subprocess.run(
                lo_cmd,
                capture_output=True,
                text=True,
                timeout=300  # 5分钟超时
            )
            
            if result.returncode != 0:
                print(f"[WARNING] LibreOffice 转换失败: {result.stderr}")
                return images
            
            if not pdf_path.exists():
                # 尝试查找生成的 PDF
                pdfs = list(temp_path.glob("*.pdf"))
                if pdfs:
                    pdf_path = pdfs[0]
                else:
                    print(f"[WARNING] PDF 文件未生成")
                    return images
            
            # 使用 pdftoppm 转换 PDF 到 PNG
            # 计算分辨率以达到 1920x1080
            # 标准幻灯片比例 16:9，如果 DPI=150，则大约 1920x1080
            ppm_cmd = [
                "pdftoppm",
                "-png",
                "-r", str(dpi),
                str(pdf_path),
                str(images_dir / pptx_path.stem)
            ]
            
            result = subprocess.run(
                ppm_cmd,
                capture_output=True,
                text=True,
                timeout=300
            )
            
            if result.returncode != 0:
                print(f"[WARNING] pdftoppm 转换失败: {result.stderr}")
                return images
            
            # 收集生成的图片
            for img in sorted(images_dir.glob(f"{pptx_path.stem}*.png")):
                images.append(str(img))
            
        except subprocess.TimeoutExpired:
            print(f"[WARNING] 转换超时")
        except FileNotFoundError as e:
            print(f"[WARNING] 命令未找到: {e}")
        except Exception as e:
            print(f"[WARNING] 转换异常: {e}")
    
    return images


def extract_pptx_content(input_path: str, output_dir: str, dpi: int = 150) -> dict[str, Any]:
    """抽取 PPTX 内容并生成图片
    
    Args:
        input_path: PPTX 文件路径或目录
        output_dir: 输出目录
        dpi: 图片 DPI（默认 150，对应约 1920x1080）
    
    Returns:
        抽取的内容（JSON 格式）
    """
    try:
        from pptx import Presentation
    except ImportError:
        return {"ok": False, "error": "python-pptx 未安装，请运行: pip install python-pptx"}
    
    input_p = Path(input_path).expanduser().resolve()
    output_p = Path(output_dir).expanduser().resolve()
    _ensure_dir(output_p)
    
    # 确定要处理的 PPTX 文件
    pptx_files = []
    if input_p.is_file():
        if input_p.suffix.lower() in ('.pptx', '.ppt'):
            pptx_files.append(input_p)
    elif input_p.is_dir():
        pptx_files.extend(sorted(input_p.glob("**/*.pptx")))
        pptx_files.extend(sorted(input_p.glob("**/*.ppt")))
    
    if not pptx_files:
        return {"ok": False, "error": f"未找到 PPTX 文件: {input_path}"}
    
    all_results = []
    
    for pptx_file in pptx_files:
        pptx_stem = pptx_file.stem
        pptx_output_dir = _ensure_dir(output_p / _safe_filename(pptx_stem))
        
        result = {
            "pptx_path": str(pptx_file),
            "pptx_stem": pptx_stem,
            "output_dir": str(pptx_output_dir),
            "pages": [],
            "images": []
        }
        
        try:
            # 解析 PPTX
            prs = Presentation(str(pptx_file))
            
            # 提取每页内容
            for idx, slide in enumerate(prs.slides):
                page_content = _extract_slide_content(slide, idx)
                result["pages"].append(page_content)
            
            # 生成图片
            images = _convert_pptx_to_images(pptx_file, pptx_output_dir, dpi)
            result["images"] = images
            
            # 将图片路径关联到页面
            for idx, page in enumerate(result["pages"]):
                if idx < len(images):
                    page["image_path"] = images[idx]
            
            # 保存 JSON 内容
            content_json_path = pptx_output_dir / "content.json"
            content_json_path.write_text(
                json.dumps(result, ensure_ascii=False, indent=2),
                encoding="utf-8"
            )
            
            # 保存 Markdown 内容
            md_lines = [f"# {pptx_stem}\n"]
            for page in result["pages"]:
                md_lines.append(f"\n## 第 {page['page_number']} 页")
                if page["title"]:
                    md_lines.append(f"\n### {page['title']}")
                for text in page["texts"]:
                    md_lines.append(f"\n{text}")
                if page["notes"]:
                    md_lines.append(f"\n**备注:** {page['notes']}")
            
            content_md_path = pptx_output_dir / "content.md"
            content_md_path.write_text("\n".join(md_lines), encoding="utf-8")
            
            result["content_json"] = str(content_json_path)
            result["content_md"] = str(content_md_path)
            
        except Exception as e:
            result["error"] = str(e)
        
        all_results.append(result)
    
    # 保存汇总结果
    summary_path = output_p / "extracted_content.json"
    summary_path.write_text(
        json.dumps({"ok": True, "results": all_results}, ensure_ascii=False, indent=2),
        encoding="utf-8"
    )
    
    return {"ok": True, "results": all_results, "summary_file": str(summary_path)}


def save_scripts(output_dir: str, scripts_json: str) -> dict[str, Any]:
    """保存讲稿到文件
    
    Args:
        output_dir: 输出目录
        scripts_json: 讲稿 JSON 字符串
    
    Returns:
        保存结果
    """
    output_p = Path(output_dir).expanduser().resolve()
    scripts_dir = _ensure_dir(output_p / "scripts")
    
    try:
        # 解析 JSON
        scripts_data = json.loads(scripts_json) if isinstance(scripts_json, str) else scripts_json
        
        # 处理可能的 JSON 外层包装
        if isinstance(scripts_data, str):
            scripts_data = json.loads(scripts_data)
        
        title = scripts_data.get("title", "PPT讲稿")
        scripts = scripts_data.get("scripts", [])
        
        saved_files = []
        
        # 保存完整 JSON
        full_json_path = scripts_dir / "scripts.json"
        full_json_path.write_text(
            json.dumps(scripts_data, ensure_ascii=False, indent=2),
            encoding="utf-8"
        )
        saved_files.append(str(full_json_path))
        
        # 保存每页讲稿为单独文件（便于编辑）
        for script in scripts:
            page_num = script.get("page_number", 0)
            script_text = script.get("script_text", "")
            duration = script.get("estimated_duration_seconds", 0)
            
            page_file = scripts_dir / f"page_{page_num:03d}.txt"
            page_file.write_text(
                f"# 第 {page_num} 页\n# 预计时长: {duration} 秒\n\n{script_text}",
                encoding="utf-8"
            )
            saved_files.append(str(page_file))
        
        # 保存合并的 Markdown 文件
        md_lines = [f"# {title}\n"]
        for script in scripts:
            page_num = script.get("page_number", 0)
            script_text = script.get("script_text", "")
            duration = script.get("estimated_duration_seconds", 0)
            md_lines.append(f"\n## 第 {page_num} 页 (约 {duration} 秒)\n")
            md_lines.append(script_text)
        
        md_path = scripts_dir / "scripts.md"
        md_path.write_text("\n".join(md_lines), encoding="utf-8")
        saved_files.append(str(md_path))
        
        return {
            "ok": True,
            "scripts_dir": str(scripts_dir),
            "saved_files": saved_files,
            "page_count": len(scripts)
        }
        
    except Exception as e:
        return {"ok": False, "error": str(e)}


def _get_speech_config():
    """获取 Azure Speech 配置"""
    try:
        import azure.cognitiveservices.speech as speechsdk
    except ImportError:
        raise RuntimeError("Azure Speech SDK 未安装，请运行: pip install azure-cognitiveservices-speech")
    
    speech_key = os.environ.get("AZURE_SPEECH_KEY")
    speech_region = os.environ.get("AZURE_SPEECH_REGION")
    speech_endpoint = os.environ.get("AZURE_SPEECH_ENDPOINT")
    
    if not speech_key:
        raise ValueError("AZURE_SPEECH_KEY 环境变量未设置")
    
    if speech_endpoint:
        return speechsdk.SpeechConfig(subscription=speech_key, endpoint=speech_endpoint)
    
    if not speech_region:
        speech_region = os.environ.get("AZURE_LOCATION")
    if speech_region:
        return speechsdk.SpeechConfig(subscription=speech_key, region=speech_region)
    
    raise ValueError("缺少 Speech 配置，请设置 AZURE_SPEECH_ENDPOINT 或 AZURE_SPEECH_REGION")


def _generate_single_audio(
    text: str,
    output_file: str,
    voice_name: str,
    speaking_rate: str,
    language: str,
    max_retries: int = 3
) -> dict[str, Any]:
    """生成单个音频文件（带重试）"""
    import azure.cognitiveservices.speech as speechsdk
    
    for attempt in range(max_retries):
        try:
            speech_config = _get_speech_config()
            speech_config.speech_synthesis_voice_name = voice_name
            speech_config.set_speech_synthesis_output_format(
                speechsdk.SpeechSynthesisOutputFormat.Audio16Khz32KBitRateMonoMp3
            )
            
            # 构建 SSML
            safe_text = _xml_escape(text)
            ssml = f'''<speak version="1.0" xmlns="http://www.w3.org/2001/10/synthesis" xml:lang="{language}">
                <voice name="{voice_name}">
                    <prosody rate="{speaking_rate}">{safe_text}</prosody>
                </voice>
            </speak>'''
            
            audio_config = speechsdk.audio.AudioOutputConfig(filename=output_file)
            synthesizer = speechsdk.SpeechSynthesizer(
                speech_config=speech_config,
                audio_config=audio_config
            )
            
            result = synthesizer.speak_ssml_async(ssml).get()
            
            if result.reason == speechsdk.ResultReason.SynthesizingAudioCompleted:
                return {"ok": True, "output_file": output_file}
            
            if result.reason == speechsdk.ResultReason.Canceled:
                cancellation = result.cancellation_details
                error_msg = f"TTS 取消: {cancellation.reason}"
                if cancellation.reason == speechsdk.CancellationReason.Error:
                    error_msg += f" - {cancellation.error_details}"
                
                if attempt < max_retries - 1:
                    wait_time = 2 ** attempt  # 指数退避
                    print(f"[TTS] 重试 {attempt + 1}/{max_retries}，等待 {wait_time} 秒...")
                    time.sleep(wait_time)
                    continue
                
                return {"ok": False, "error": error_msg}
            
        except Exception as e:
            if attempt < max_retries - 1:
                wait_time = 2 ** attempt
                print(f"[TTS] 异常重试 {attempt + 1}/{max_retries}，等待 {wait_time} 秒: {e}")
                time.sleep(wait_time)
                continue
            return {"ok": False, "error": str(e)}
    
    return {"ok": False, "error": "TTS 生成失败（已达最大重试次数）"}


def generate_page_audio(
    output_dir: str,
    scripts_json: str,
    voice_name: str = "zh-CN-XiaoxiaoNeural",
    speaking_rate: str = "1.0",
    language: str = "zh-CN"
) -> dict[str, Any]:
    """为每一页生成 TTS 音频
    
    Args:
        output_dir: 输出目录
        scripts_json: 讲稿 JSON
        voice_name: 音色名称
        speaking_rate: 语速
        language: 语言
    
    Returns:
        生成的音频文件列表
    """
    output_p = Path(output_dir).expanduser().resolve()
    audio_dir = _ensure_dir(output_p / "audio")
    
    try:
        # 解析 JSON
        scripts_data = json.loads(scripts_json) if isinstance(scripts_json, str) else scripts_json
        if isinstance(scripts_data, str):
            scripts_data = json.loads(scripts_data)
        
        scripts = scripts_data.get("scripts", [])
        if not scripts:
            return {"ok": False, "error": "讲稿为空"}
        
        audio_files = []
        errors = []
        
        # 串行生成（避免 API 限流）
        for script in scripts:
            page_num = script.get("page_number", 0)
            script_text = script.get("script_text", "")
            
            if not script_text.strip():
                continue
            
            output_file = str(audio_dir / f"page_{page_num:03d}.mp3")
            result = _generate_single_audio(
                text=script_text,
                output_file=output_file,
                voice_name=voice_name,
                speaking_rate=speaking_rate,
                language=language
            )
            
            if result.get("ok"):
                audio_files.append({
                    "page_number": page_num,
                    "audio_file": output_file
                })
                print(f"[TTS] 第 {page_num} 页音频生成完成: {output_file}")
            else:
                errors.append({
                    "page_number": page_num,
                    "error": result.get("error", "未知错误")
                })
        
        # 保存结果
        result_path = audio_dir / "audio_files.json"
        result_path.write_text(
            json.dumps({"audio_files": audio_files, "errors": errors}, ensure_ascii=False, indent=2),
            encoding="utf-8"
        )
        
        return {
            "ok": len(errors) == 0,
            "audio_files": audio_files,
            "errors": errors,
            "audio_dir": str(audio_dir)
        }
        
    except Exception as e:
        return {"ok": False, "error": str(e)}


def _get_audio_duration(audio_file: str) -> float:
    """获取音频时长（秒）"""
    try:
        result = subprocess.run(
            ["ffprobe", "-v", "quiet", "-show_entries", "format=duration",
             "-of", "default=noprint_wrappers=1:nokey=1", audio_file],
            capture_output=True,
            text=True,
            timeout=30
        )
        return float(result.stdout.strip())
    except Exception:
        return 0.0


def _generate_single_video(
    image_file: str,
    audio_file: str,
    output_file: str
) -> dict[str, Any]:
    """生成单个视频片段"""
    try:
        duration = _get_audio_duration(audio_file)
        if duration <= 0:
            return {"ok": False, "error": f"无法获取音频时长: {audio_file}"}
        
        # 使用 ffmpeg 生成视频
        # -loop 1: 循环图片
        # -i: 输入图片和音频
        # -c:v libx264: 视频编码
        # -tune stillimage: 针对静态图片优化
        # -c:a aac: 音频编码
        # -b:a 192k: 音频比特率
        # -pix_fmt yuv420p: 像素格式（兼容性）
        # -shortest: 以最短流为准
        cmd = [
            "ffmpeg", "-y",
            "-loop", "1",
            "-i", image_file,
            "-i", audio_file,
            "-c:v", "libx264",
            "-tune", "stillimage",
            "-c:a", "aac",
            "-b:a", "192k",
            "-pix_fmt", "yuv420p",
            "-shortest",
            output_file
        ]
        
        result = subprocess.run(
            cmd,
            capture_output=True,
            text=True,
            timeout=300
        )
        
        if result.returncode != 0:
            return {"ok": False, "error": f"ffmpeg 错误: {result.stderr}"}
        
        return {"ok": True, "output_file": output_file, "duration": duration}
        
    except Exception as e:
        return {"ok": False, "error": str(e)}


def generate_video(output_dir: str) -> dict[str, Any]:
    """生成视频片段并合成最终视频
    
    Args:
        output_dir: 输出目录（包含 images/ 和 audio/ 子目录）
    
    Returns:
        最终视频路径
    """
    output_p = Path(output_dir).expanduser().resolve()
    video_dir = _ensure_dir(output_p / "video")
    clips_dir = _ensure_dir(video_dir / "clips")
    
    # 查找图片和音频文件
    images_dir = output_p / "images"
    audio_dir = output_p / "audio"
    
    # 如果 images 在子目录中，尝试查找
    if not images_dir.exists():
        for subdir in output_p.iterdir():
            if subdir.is_dir():
                candidate = subdir / "images"
                if candidate.exists():
                    images_dir = candidate
                    break
    
    if not images_dir.exists():
        return {"ok": False, "error": f"图片目录不存在: {images_dir}"}
    
    if not audio_dir.exists():
        return {"ok": False, "error": f"音频目录不存在: {audio_dir}"}
    
    # 收集图片和音频（按页码排序）
    images = sorted(images_dir.glob("*.png"))
    audios = sorted(audio_dir.glob("*.mp3"))
    
    if not images:
        return {"ok": False, "error": "未找到图片文件"}
    
    if not audios:
        return {"ok": False, "error": "未找到音频文件"}
    
    # 生成视频片段（并行）
    clips = []
    errors = []
    
    with ThreadPoolExecutor(max_workers=4) as executor:
        futures = {}
        
        for idx, (img, audio) in enumerate(zip(images, audios)):
            clip_file = str(clips_dir / f"clip_{idx + 1:03d}.mp4")
            future = executor.submit(
                _generate_single_video,
                str(img),
                str(audio),
                clip_file
            )
            futures[future] = (idx + 1, clip_file)
        
        for future in as_completed(futures):
            page_num, clip_file = futures[future]
            result = future.result()
            
            if result.get("ok"):
                clips.append({
                    "page_number": page_num,
                    "clip_file": clip_file,
                    "duration": result.get("duration", 0)
                })
                print(f"[Video] 第 {page_num} 页视频片段生成完成")
            else:
                errors.append({
                    "page_number": page_num,
                    "error": result.get("error", "未知错误")
                })
    
    if not clips:
        return {"ok": False, "error": "没有成功生成任何视频片段", "errors": errors}
    
    # 按页码排序
    clips.sort(key=lambda x: x["page_number"])
    
    # 使用 ffmpeg concat 合成最终视频（无重编码）
    concat_list_file = video_dir / "concat_list.txt"
    with open(concat_list_file, "w", encoding="utf-8") as f:
        for clip in clips:
            f.write(f"file '{clip['clip_file']}'\n")
    
    # 获取 PPTX 名称作为输出文件名
    pptx_stem = "output"
    for subdir in output_p.iterdir():
        if subdir.is_dir() and subdir.name not in ("audio", "video", "scripts", "images"):
            pptx_stem = subdir.name
            break
    
    final_video = video_dir / f"{pptx_stem}.mp4"
    
    try:
        # 使用 concat demuxer 无重编码拼接
        cmd = [
            "ffmpeg", "-y",
            "-f", "concat",
            "-safe", "0",
            "-i", str(concat_list_file),
            "-c", "copy",
            str(final_video)
        ]
        
        result = subprocess.run(
            cmd,
            capture_output=True,
            text=True,
            timeout=600
        )
        
        if result.returncode != 0:
            return {"ok": False, "error": f"视频合成失败: {result.stderr}"}
        
        total_duration = sum(c.get("duration", 0) for c in clips)
        
        return {
            "ok": True,
            "final_video": str(final_video),
            "total_duration_seconds": total_duration,
            "clip_count": len(clips),
            "errors": errors
        }
        
    except Exception as e:
        return {"ok": False, "error": str(e)}


def tavily_search_for_ppt(content_summary: str, max_searches: int = 5) -> dict[str, Any]:
    """使用 Tavily 搜索 PPT 相关内容
    
    Args:
        content_summary: PPT 内容摘要
        max_searches: 最大搜索次数
    
    Returns:
        搜索结果摘要
    """
    try:
        from tavily import TavilyClient
    except ImportError:
        return {"ok": False, "error": "Tavily 客户端未安装，请运行: pip install tavily-python"}
    
    api_key = os.environ.get("TAVILY_API_KEY")
    if not api_key:
        return {"ok": False, "error": "TAVILY_API_KEY 环境变量未设置"}
    
    client = TavilyClient(api_key=api_key)
    
    # 从内容摘要中提取关键词作为搜索查询
    # 简单实现：取前 500 字符作为上下文
    context = (content_summary or "")[:500]
    
    try:
        # 执行搜索
        results = []
        
        # 生成搜索查询（基于内容）
        queries = [
            f"{context[:100]} 最新进展",
            f"{context[:100]} 案例分析",
            f"{context[:100]} 最佳实践"
        ][:max_searches]
        
        for query in queries:
            try:
                resp = client.search(
                    query=query,
                    max_results=3,
                    include_answer=True,
                    search_depth="basic"
                )
                results.append({
                    "query": query,
                    "answer": resp.get("answer", ""),
                    "sources": [
                        {"title": r.get("title", ""), "url": r.get("url", "")}
                        for r in resp.get("results", [])[:3]
                    ]
                })
            except Exception as e:
                results.append({"query": query, "error": str(e)})
        
        return {
            "ok": True,
            "search_count": len(results),
            "results": results
        }
        
    except Exception as e:
        return {"ok": False, "error": str(e)}


def register_tools(registry: object) -> None:
    """注册所有 PPT 处理工具"""
    register = getattr(registry, "register_tool", None)
    if not callable(register):
        return
    
    register("ppt.scan_pptx_files", scan_pptx_files)
    register("ppt.extract_pptx_content", extract_pptx_content)
    register("ppt.save_scripts", save_scripts)
    register("ppt.generate_page_audio", generate_page_audio)
    register("ppt.generate_video", generate_video)
    register("ppt.tavily_search_for_ppt", tavily_search_for_ppt)


# CLI 入口
if __name__ == "__main__":
    import argparse
    
    parser = argparse.ArgumentParser(description="PPT Processor Tool")
    parser.add_argument("--tool", required=True, help="工具名称")
    parser.add_argument("--args-json", help="工具参数 (JSON)")
    
    args = parser.parse_args()
    
    tool_map = {
        "scan": scan_pptx_files,
        "extract": extract_pptx_content,
        "save_scripts": save_scripts,
        "generate_audio": generate_page_audio,
        "generate_video": generate_video,
        "tavily_search": tavily_search_for_ppt,
    }
    
    if args.tool not in tool_map:
        print(f"未知工具: {args.tool}")
        print(f"可用工具: {', '.join(tool_map.keys())}")
        exit(1)
    
    tool_args = json.loads(args.args_json) if args.args_json else {}
    result = tool_map[args.tool](**tool_args)
    print(json.dumps(result, ensure_ascii=False, indent=2))
